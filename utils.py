import logging
import os
import re
import subprocess
import unicodedata
from pathlib import Path
from io import BytesIO
from urllib.parse import urlparse

import fitz
import pandas as pd
import pytesseract
import sqlparse
from PIL import Image
from docx import Document as DocxDocument
from pptx import Presentation
from dotenv import load_dotenv
from pathspec import PathSpec
from sqlalchemy import create_engine, text
from unstructured.partition.auto import partition

load_dotenv()

def load_prompt(filepath):
    return Path(filepath).read_text(encoding="utf-8")

ANTHROPIC_API_KEY = os.getenv("ANTHROPIC_API_KEY")
CLAUDE_MODEL = os.getenv("CLAUDE_MODEL", "claude-3-haiku-20240307")

SANDBOX_CONTAINER_NAME = os.getenv("SANDBOX_CONTAINER_NAME", "rag-assistant-rag-sandbox-1")
_ENGINES = {}


REPOS_ROOT = Path("repos").resolve()
REPOS_SAFE_ROOT = Path("repos_safe").resolve()

def setup_logging(name: str, file: bool = True) -> logging.Logger:
    logger = logging.getLogger(name)
    if logger.handlers:
        return logger
    logger.setLevel(logging.INFO)
    fmt = logging.Formatter("%(asctime)s | %(levelname)s | %(message)s")
    if file:
        log_dir = Path("logs")
        log_dir.mkdir(exist_ok=True)
        fh = logging.FileHandler(log_dir / f"{name}.log", encoding="utf-8")
        fh.setFormatter(fmt)
        logger.addHandler(fh)
    ch = logging.StreamHandler()
    ch.setFormatter(fmt)
    logger.addHandler(ch)
    logger.propagate = False
    return logger

def to_posix(p: str | Path) -> str:
    s = (str(p) or "").strip().replace("\\", "/")
    while s.startswith("./") or s.startswith("../"):
        s = s[2:] if s.startswith("./") else s[3:]
    while s.startswith("/"):
        s = s[1:]
    while "//" in s:
        s = s.replace("//", "/")
    return s

def clean_text(text: str) -> str:
    text = unicodedata.normalize("NFKC", text)
    text = re.sub(r"[\u200b\u200e\u200f\u202a-\u202e]", "", text)
    text = re.sub(r"(\w)-\n(\w)", r"\1\2", text)
    text = re.sub(r"\s*\n\s*", "\n", text)
    text = re.sub(r"[ \t]+", " ", text)
    return text.strip()

IGNORE_FILE = Path(".ignore")
if not IGNORE_FILE.exists():
    raise FileNotFoundError(f"Файл .ignore не найден в {IGNORE_FILE.resolve()}")
IGNORE_SPEC = PathSpec.from_lines("gitwildmatch", [stripped for line in IGNORE_FILE.read_text(encoding="utf-8").splitlines() if (stripped := line.partition("#")[0].strip())])

def is_ignored(rel_path: str) -> bool:
    return IGNORE_SPEC.match_file(rel_path)

def extract_binary_content(path: Path):
    with open(path, "rb") as f:
        file_bytes = f.read()
    ext = path.suffix.lower()
    if ext == ".pdf":
        doc = fitz.open(stream=BytesIO(file_bytes), filetype="pdf")
        return "\n".join([page.get_text() for page in doc])
    elif ext == ".docx":
        docx_doc = DocxDocument(BytesIO(file_bytes))
        return "\n".join([para.text for para in docx_doc.paragraphs])
    elif ext == ".pptx":
        prs = Presentation(BytesIO(file_bytes))
        text_parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                text_parts.append(shape.text)
        return "\n".join(text_parts)
    elif ext in [".html", ".epub", ".rtf"]:
        elements = partition(file=BytesIO(file_bytes))
        return "\n".join([el.text for el in elements if getattr(el, "text", "").strip()])
    elif ext == ".csv":
        df = pd.read_csv(BytesIO(file_bytes))
        return df.to_string(index=False)
    elif ext in [".xls", ".xlsx"]:
        df = pd.read_excel(BytesIO(file_bytes))
        return df.to_string(index=False)
    elif ext in [".png", ".jpg", ".jpeg", ".tif", ".tiff", ".webp", ".bmp"]:
        img = Image.open(BytesIO(file_bytes))
        return pytesseract.image_to_string(img, lang="rus+eng")
    return None

def execute_command(command: str):
    ps_cmd = ['docker', 'ps', '--filter', f'name={SANDBOX_CONTAINER_NAME}', '--format', '{{.ID}}']
    ps_result = subprocess.run(ps_cmd, capture_output=True, text=True)
    if ps_result.returncode != 0 or not ps_result.stdout.strip():
        raise RuntimeError(f"Контейнер {SANDBOX_CONTAINER_NAME} не найден")
    container_id = ps_result.stdout.strip()
    exec_cmd = ['docker', 'exec', '-u', 'nobody', container_id, 'timeout', '30', 'sh', '-c', command]
    exec_result = subprocess.run(exec_cmd, capture_output=True, text=True)
    return {
        "stdout": exec_result.stdout,
        "exit_code": exec_result.returncode
    }

def load_db_connections():
    connections = {}
    idx = 1
    while True:
        url = os.getenv(f"DB_{idx}_URL")
        if not url:
            break
        username = os.getenv(f"DB_{idx}_USERNAME")
        password = os.getenv(f"DB_{idx}_PASSWORD")
        description = os.getenv(f"DB_{idx}_DESCRIPTION")
        tool_name = os.getenv(f"DB_{idx}_TOOL_NAME")
        if username and password and tool_name:
            connections[tool_name] = {
                "url": url,
                "username": username,
                "password": password,
                "description": description,
            }
        idx += 1
    return connections

DB_CONNECTIONS = load_db_connections()

def _sqlalchemy_url(conn: dict) -> str:
    raw = (conn.get("url") or "").strip().removeprefix("jdbc:")
    if "://" not in raw:
        raw = "//" + raw
    parsed = urlparse(raw)
    scheme = parsed.scheme or "postgresql"
    host = parsed.hostname or ""
    port = f":{parsed.port}" if parsed.port else ""
    db = (parsed.path or "").lstrip("/") or conn.get("db_name", "")
    return f"{scheme}://{conn['username']}:{conn['password']}@{host}{port}/{db}"

def get_engine(tool_name: str):
    if tool_name in _ENGINES:
        return _ENGINES[tool_name]
    conn = DB_CONNECTIONS[tool_name]
    url = _sqlalchemy_url(conn)
    connect_args = {}
    if url.startswith("postgresql"):
        connect_args = {"options": "-c statement_timeout=30s"}
    engine = create_engine(
        url,
        connect_args=connect_args,
        pool_pre_ping=True,
        pool_timeout=10,
        pool_recycle=3600,
    )
    _ENGINES[tool_name] = engine
    return engine

def db_query(tool_name, select_query):
    parsed = sqlparse.parse((select_query or "").strip())
    if not parsed or parsed[0].get_type() != "SELECT":
        return {"error": "Разрешены только SELECT запросы", "rows": [], "select": select_query}
    try:
        engine = get_engine(tool_name)
        with engine.connect() as connection:
            result = connection.execute(text(select_query))
            rows = [dict(row._mapping) for row in result]
        return {"rows": rows, "select": select_query}
    except Exception as exc:
        return {"error": str(exc), "rows": [], "select": select_query}

def build_select_tools():
    tools = []
    for tool_name, conn in DB_CONNECTIONS.items():
        tools.append({
            "name": tool_name,
            "description": conn["description"],
            "input_schema": {
                "type": "object",
                "properties": {
                    "purpose": {
                        "type": "string",
                        "description": "Зачем этот SELECT: опиши цель запроса."
                    },
                    "select": {
                        "type": "string",
                        "description": (
                            "Готовый SQL SELECT запрос. "
                            "Разрешены только операции чтения. "
                            "НЕ передавай сюда вопросы на естественном языке. "
                            "Только валидный SELECT."
                        )
                    }
                },
                "required": ["purpose", "select"]
            }
        })
    return tools

SELECT_TOOLS = build_select_tools()

EXECUTE_COMMAND_TOOL = {
    "name": "execute_command",
    "description": "Выполнение любых консольных команд в изолированном контейнере для взаимодействия с исходными файлами.",
    "input_schema": {
        "type": "object",
        "properties": {
            "purpose": {
                "type": "string",
                "description": "Зачем запускается команда: что хотим узнать."
            },
            "command": {
                "type": "string",
                "description": (
                    "команда для выполнения в изолированном контейнере. "
                    "Доступные утилиты: grep, find, awk, sed, bash, git, jq, tree, file, diff, less, vim, cat, head, tail, wc, sort, uniq, cut, tr, xargs, "
                    "basename, dirname, realpath, stat, ls, cmp, split, tee, seq, od, strings, tar, gzip, md5sum, sha256sum, date, env, ps, df, du, "
                    "which, type, command, test, readlink, echo, printf"
                )
            }
        },
        "required": ["purpose", "command"]
    }
}

