"""Sourcegraph инструменты для работы с кодом через GraphQL API"""

import requests
from pathlib import Path
from utils import SOURCEGRAPH_URL, SOURCEGRAPH_TOKEN, setup_logging, clean_text, REPOS_ROOT
from mask import mask_secrets, check_secrets_in_text

logger = setup_logging(Path(__file__).stem)

GRAPHQL_ENDPOINT = f"{SOURCEGRAPH_URL}/.api/graphql"
HEADERS = {
    "Content-Type": "application/json",
    "Authorization": f"token {SOURCEGRAPH_TOKEN}" if SOURCEGRAPH_TOKEN else "",
}

BINARY_CHUNK_SIZE = 512

import pytesseract
from PIL import Image
from io import BytesIO
import fitz
import pandas as pd
from docx import Document as DocxDocument
from pptx import Presentation
from llama_index.readers.file import UnstructuredReader
from llama_index.core import Document

def _execute_graphql(query: str, variables: dict | None = None) -> dict:
    payload = {"query": query}
    if variables:
        payload["variables"] = variables
    resp = requests.post(GRAPHQL_ENDPOINT, json=payload, headers=HEADERS, timeout=30)
    resp.raise_for_status()
    return resp.json()

def _split_file_path(rel_path: str) -> tuple[str, str]:
    parts = rel_path.split("/", 1)
    if len(parts) < 2:
        raise ValueError(f"Invalid file path format: {rel_path}. Expected format: repo/path")
    return parts[0], parts[1]

def _extract_binary_content(rel_path: str, branch: str):
    repo, path = _split_file_path(rel_path)
    raw_url = f"{SOURCEGRAPH_URL}/{repo}@{branch}/-/raw/{path}"
    resp = requests.get(raw_url, headers=HEADERS, timeout=30)
    resp.raise_for_status()
    file_bytes = resp.content
    ext = Path(rel_path).suffix.lower()
    try:
        if ext == ".pdf":
            doc = fitz.open(stream=BytesIO(file_bytes), filetype="pdf")
            text = "\n".join([page.get_text() for page in doc])
            docs = [Document(text=text, metadata={"file_path": rel_path})]
        elif ext == ".docx":
            docx_doc = DocxDocument(BytesIO(file_bytes))
            text = "\n".join([para.text for para in docx_doc.paragraphs])
            docs = [Document(text=text, metadata={"file_path": rel_path})]
        elif ext == ".pptx":
            prs = Presentation(BytesIO(file_bytes))
            text_parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_parts.append(shape.text)
            docs = [Document(text="\n".join(text_parts), metadata={"file_path": rel_path})]
        elif ext in [".html", ".epub", ".rtf"]:
            reader = UnstructuredReader()
            docs = reader.load_data(file=BytesIO(file_bytes))
        elif ext == ".csv":
            df = pd.read_csv(BytesIO(file_bytes))
            text = df.to_string(index=False)
            docs = [Document(text=text, metadata={"file_path": rel_path})]
        elif ext in [".xls", ".xlsx"]:
            df = pd.read_excel(BytesIO(file_bytes))
            text = df.to_string(index=False)
            docs = [Document(text=text, metadata={"file_path": rel_path})]
        elif ext in [".png", ".jpg", ".jpeg", ".tif", ".tiff", ".webp", ".bmp"]:
            img = Image.open(BytesIO(file_bytes))
            try:
                txt = pytesseract.image_to_string(img, lang="rus+eng")
            except Exception:
                txt = pytesseract.image_to_string(img)
            docs = [Document(text=txt or "", metadata={"file_path": rel_path})]
        else:
            return None
        if not docs or not docs[0].text:
            return None
        return "\n".join(doc.text for doc in docs if doc.text)
    except Exception as e:
        logger.warning(f"Failed to extract binary content from {rel_path}: {e}")
        return None

def _create_chunk(start_line: int, end_line: int, title: str, kind: str, text: str = None) -> dict:
    chunk = {
        "start_line": start_line,
        "end_line": end_line,
        "title": title,
        "kind": kind,
    }
    if text is not None:
        chunk["text"] = text
    return chunk

def _split_binary_content(content: str) -> list[dict]:
    lines = content.split("\n")
    chunks = []
    current_start = 1
    current_length = 0
    for line_num, line in enumerate(lines, start=1):
        line_length = len(line) + 1
        if current_length + line_length > BINARY_CHUNK_SIZE and current_start < line_num:
            chunks.append(_create_chunk(current_start, line_num - 1, "binary_content", "binary"))
            current_start = line_num
            current_length = line_length
        else:
            current_length += line_length
    if current_start <= len(lines):
        chunks.append(_create_chunk(current_start, len(lines), "binary_content", "binary"))
    return chunks

def sanitize_chunk(text: str) -> str:
    masked = mask_secrets(text)
    check_secrets_in_text(masked)
    return masked

def get_file_chunks(rel_path: str, branch: str) -> list[dict]:
    repo, path = _split_file_path(rel_path)
    gql = """
    query FileChunks($repo: String!, $path: String!, $rev: String!) {
      repository(name: $repo) {
        commit(rev: $rev) {
          file(path: $path) {
            content
            binary
            symbols(first: 2000) {
              nodes {
                name
                kind
                range { start { line } end { line } }
              }
            }
          }
        }
      }
    }
    """
    result = _execute_graphql(gql, {"repo": repo, "path": path, "rev": branch})
    file_data = result["data"]["repository"]["commit"]["file"]
    if file_data.get("binary"):
        content = _extract_binary_content(rel_path, branch)
        if content is None:
            return []
        content = clean_text(content)
        chunks_meta = _split_binary_content(content)
    else:
        content = file_data.get("content")
        if content is None:
            return []
        symbols = (file_data.get("symbols", {}) or {}).get("nodes", [])
        if symbols:
            chunks_meta = []
            for s in symbols:
                r = s.get("range", {})
                start = max(1, (r.get("start", {}) or {}).get("line", 1))
                end = max(start, (r.get("end", {}) or {}).get("line", start))
                chunks_meta.append(_create_chunk(start, end, s.get("name", ""), s.get("kind", "symbol")))
            chunks_meta.sort(key=lambda x: (x["start_line"], x["end_line"]))
        else:
            content = clean_text(content)
            chunks_meta = _split_binary_content(content)
    lines = content.split("\n")
    chunks = []
    for ch in chunks_meta:
        start = max(1, int(ch["start_line"]))
        end = max(start, int(ch["end_line"]))
        text = "\n".join(lines[start-1:end])
        if text:
            text = sanitize_chunk(text)
            if text:
                chunks.append(_create_chunk(start, end, ch["title"], ch["kind"], text))
    return chunks

def get_file_hash(rel_path: str, branch: str) -> str | None:
    repo, path = _split_file_path(rel_path)
    gql_query = """
    query BlobOid($repo: String!, $path: String!, $rev: String!) {
      repository(name: $repo) { commit(rev: $rev) { file(path: $path) { oid } } }
    }
    """
    result = _execute_graphql(gql_query, {"repo": repo, "path": path, "rev": branch})
    file_data = result["data"]["repository"]["commit"]["file"]
    return file_data.get("oid")

def sg_search(query: str, repo: str, limit: int) -> str:
    repo_filter = f"repo:{repo}" if repo else ""
    search_query = f"{repo_filter} {query}" if repo_filter else query
    gql_query = """
    query SearchResults($query: String!, $limit: Int!) {
      search(query: $query, first: $limit) {
        results { results { ... on FileMatch { file { path url } lineMatches { lineNumber preview } } } }
      }
    }
    """
    result = _execute_graphql(gql_query, {"query": search_query, "limit": limit})
    matches = result["data"]["search"]["results"]["results"]
    if not matches:
        return f"Поиск: '{query}' не дал результатов"
    out = [f"🔍 Sourcegraph поиск: '{query}' ({len(matches)} результатов):\n"]
    for m in matches:
        f = m["file"]
        out.append(f"\n📄 {f['path']}")
        for lm in m.get("lineMatches", [])[:5]:
            out.append(f"  {lm.get('lineNumber','')}: {lm.get('preview','').strip()}")
    return "\n".join(out)

def sg_codeintel(mode: str, symbol: str, doc_id: str, line: int) -> str:
    gql_query_map = {
        "definitions": """
          query Definitions($symbol: String!) {
            symbolSearch(query: $symbol, first: 10) {
              results { result { file { path } symbol { name containerName } locations { range { start { line } end { line } } } } }
            }
          }
        """,
        "references": """
          query References($symbol: String!) {
            symbolReferences(query: $symbol, first: 20) {
              references { file { path } symbol { name } location { range { start { line } end { line } } } }
            }
          }
        """,
    }
    gql_query = gql_query_map.get(mode)
    if not gql_query:
        return f"Неизвестный режим: {mode}. Доступные: definitions, references, callers, callees"
    data = _execute_graphql(gql_query, {"symbol": symbol})["data"]
    if mode == "definitions":
        defs = data["symbolSearch"]["results"]
        if not defs:
            return f"Определения для '{symbol}' не найдены"
        out = [f"📌 Определения символа '{symbol}':\n"]
        for d in defs:
            r = d["result"]
            out.append(f"  📄 {r['file']['path']}")
            out.append(f"     Символ: {r['symbol'].get('name','')}")
            for loc in r.get("locations", [])[:3]:
                out.append(f"     Строки: {loc['range']['start'].get('line','N/A')}")
        return "\n".join(out)
    if mode == "references":
        refs = data["symbolReferences"]["references"]
        if not refs:
            return f"Ссылки на '{symbol}' не найдены"
        out = [f"🔗 Ссылки на символ '{symbol}' ({len(refs)} найден):\n"]
        for r in refs[:15]:
            out.append(f"  📄 {r['file'].get('path','unknown')}:{r['location']['range']['start'].get('line','N/A')}")
        return "\n".join(out)
    return f"Обработан режим: {mode}"

def sg_blob(doc_id: str, start_line: int, end_line: int, branch: str) -> str:
    if not doc_id:
        return "Необходимо указать doc_id"
    gql_query = """
    query BlobContent($repo: String!, $path: String!, $rev: String!) {
      repository(name: $repo) { commit(rev: $rev) { file(path: $path) { content binary } } }
    }
    """
    repo, file_path = doc_id.split("/", 1)
    result = _execute_graphql(gql_query, {"repo": repo, "path": file_path, "rev": branch})
    file_data = result["data"]["repository"]["commit"]["file"]
    if not file_data or file_data.get("binary"):
        return f"Файл {doc_id} не найден" if not file_data else f"Файл {doc_id} является бинарным"
    lines = file_data.get("content", "").split("\n")
    end_line = min(end_line, len(lines))
    start_line = max(1, start_line)
    selected = lines[start_line - 1:end_line]
    out = [f"📄 {doc_id} (строки {start_line}-{end_line}):\n"]
    for i, line in enumerate(selected, start=start_line):
        out.append(f"{i:4d} | {line}")
    return "\n".join(out)


def sg_list_repos(prefix: str = "") -> list[str]:
    gql = """
    query Repos($first: Int!) { repositories(first: $first) { nodes { name } } }
    """
    data = _execute_graphql(gql, {"first": 500})["data"]["repositories"]["nodes"]
    names = [n["name"] for n in data if n.get("name")]
    return [n for n in names if n.startswith(prefix)] if prefix else names

def sg_list_repo_branches(repo: str) -> list[str]:
    gql = """
    query RepoBranches($repo: String!) {
      repository(name: $repo) {
        branches(first: 100) {
          nodes { name }
        }
      }
    }
    """
    data = _execute_graphql(gql, {"repo": repo})["data"]["repository"]
    branches = data.get("branches", {}).get("nodes", [])
    return [b["name"] for b in branches if b.get("name")]

def sg_get_repo_branch(repo: str) -> str:
    repo_path = REPOS_ROOT / repo
    if not repo_path.exists():
        raise ValueError(f"Repository {repo} not found in {REPOS_ROOT}")
    git_head = repo_path / ".git" / "HEAD"
    if not git_head.exists():
        raise ValueError(f"Repository {repo} is not a git repository")
    head_content = git_head.read_text().strip()
    if head_content.startswith("ref: refs/heads/"):
        return head_content.replace("ref: refs/heads/", "")
    raise ValueError(f"Repository {repo} is in detached HEAD state")

def get_all_repos_branches_formatted() -> str | None:
    repos = sg_list_repos(prefix="")
    if not repos:
        return ""
    lines = ["## Доступные ветки репозиториев:\n"]
    for repo in sorted(repos):
        branch_names = sg_list_repo_branches(repo)
        if branch_names:
            branches_str = ", ".join(branch_names[:10])
            if len(branch_names) > 10:
                branches_str += f" (+{len(branch_names)-10} еще)"
            lines.append(f"- **{repo}**: {branches_str}")
    if len(lines) == 1:
        return ""
    return "\n".join(lines) + "\n"

def sg_list_repo_files(repo: str, branch: str, limit: int = 5000) -> list[tuple[str, str | None]]:
    gql = """
    query RepoFiles($query: String!, $limit: Int!) {
      search(query: $query, first: $limit) { results { results { ... on FileMatch { file { path } } } } }
    }
    """
    q = f"repo:{repo}@{branch} type:file"
    results = _execute_graphql(gql, {"query": q, "limit": limit})["data"]["search"]["results"]["results"]
    files = []
    for r in results:
        if r.get("file"):
            file_path = r["file"]["path"]
            file_rel_path = f"{repo}/{file_path}"
            oid = get_file_hash(file_rel_path, branch)
            files.append((file_rel_path, oid))
    return files

